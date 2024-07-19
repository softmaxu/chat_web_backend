import os
import sys
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

class DocumentParser:
    def __init__(self, file_path):
        self.file_path = file_path
        self.file_extension = os.path.splitext(file_path)[1].lower()

    def extract_text(self):
        if self.file_extension == '.txt':
            return self.extract_text_from_txt()
        elif self.file_extension == '.pdf':
            return self.extract_text_from_pdf()
        elif self.file_extension == '.docx':
            return self.extract_text_from_docx()
        elif self.file_extension == '.pptx':
            return self.extract_text_from_pptx()
        else:
            raise ValueError(f"Unsupported file extension: {self.file_extension}")

    def extract_text_from_txt(self):
        with open(self.file_path, 'r', encoding='utf-8') as file:
            return file.read()

    def extract_text_from_pdf(self):
        reader = PdfReader(self.file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        return text

    def extract_text_from_docx(self):
        doc = Document(self.file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + '\n'
        return text

    def extract_text_from_pptx(self):
        prs = Presentation(self.file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text
    
    def _write_text_to_file(self, text, output_file=None):
        if not output_file:
            output_file = os.path.splitext(self.file_path)[0] + '.txt'
        with open(output_file, 'w', encoding='utf-8') as file:
            file.write(text)
        return output_file
    
    def extract_text(self):
        try:
            if self.file_extension.lower() == '.txt':
                return {"ok": True,"file":self._write_text_to_file(self.extract_text_from_txt()),"msg":""}
            elif self.file_extension.lower() == '.pdf':
                return {"ok": True,"file":self._write_text_to_file(self.extract_text_from_pdf()),"msg":""}
            elif self.file_extension.lower() == '.docx':
                return {"ok": True,"file":self._write_text_to_file(self.extract_text_from_docx()),"msg":""}
            elif self.file_extension.lower() == '.pptx':
                return {"ok": True,"file":self._write_text_to_file(self.extract_text_from_pptx()),"msg":""}
            else:
                raise ValueError(f"Unsupported file extension: {self.file_extension}")
        except Exception as e:
            return {"ok": False,"file":self.file_path,"msg":str(e)}

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python extract_text.py <file_path>")
        sys.exit(1)

    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        sys.exit(1)

    try:
        parser = DocumentParser(file_path)
        text = parser.extract_text()
        print(f"Extracted text from {file_path}:\n")
        print(text)
    except ValueError as e:
        print(e)
        sys.exit(1)
