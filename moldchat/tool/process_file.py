import os
import sys
import re
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import opencc

class DocumentParser:
    def __init__(self, file_path):
        self.file_path = file_path
        self.file_extension = os.path.splitext(file_path)[1].lower()

    def extract_text(self):
        if self.file_extension == '.txt':
            return self.extract_text_from_txt()
        elif self.file_extension == '.pdf':
            return self.extract_text_from_pdf()
        elif self.file_extension == '.docx':
            return self.extract_text_from_docx()
        elif self.file_extension == '.pptx':
            return self.extract_text_from_pptx()
        else:
            raise ValueError(f"Unsupported file extension: {self.file_extension}")

    def extract_text_from_txt(self):
        with open(self.file_path, 'r', encoding='utf-8') as file:
            return file.read()

    def extract_text_from_pdf(self):
        reader = PdfReader(self.file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        return text

    def extract_text_from_docx(self):
        doc = Document(self.file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + '\n'
        return text

    def extract_text_from_pptx(self):
        prs = Presentation(self.file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text
    
    def _write_text_to_file(self, text, output_file=None):
        if not output_file:
            output_file = os.path.splitext(self.file_path)[0] + '.txt'
        with open(output_file, 'w', encoding='utf-8') as file:
            file.write(text)
        return output_file
    
    def extract_text(self):
        try:
            if self.file_extension.lower() == '.txt':
                return {"ok": True,"file":self._write_text_to_file(self.extract_text_from_txt()),"msg":""}
            elif self.file_extension.lower() == '.pdf':
                return {"ok": True,"file":self._write_text_to_file(self.extract_text_from_pdf()),"msg":""}
            elif self.file_extension.lower() == '.docx':
                return {"ok": True,"file":self._write_text_to_file(self.extract_text_from_docx()),"msg":""}
            elif self.file_extension.lower() == '.pptx':
                return {"ok": True,"file":self._write_text_to_file(self.extract_text_from_pptx()),"msg":""}
            else:
                raise ValueError(f"Unsupported file extension: {self.file_extension}")
        except Exception as e:
            return {"ok": False,"file":self.file_path,"msg":str(e)}
        
class TextCleaner:
    def clean(self, in_file_path : str, out_file_path : str):
        with open(in_file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        text = self.remove_spaces(text)
        text = self.zh_traditional_to_simplified(text)
        with open(out_file_path, 'w', encoding='utf-8') as file:
            file.write(text)
        return text
    
    def remove_spaces(self, text):
        # 更新正则表达式以覆盖更广泛的情况，包括中文字符、中文标点和部分特殊符号之间的空格
        # \u4e00-\u9fff 是常用汉字范围，\u3000-\u303f 包括了部分中文标点和特殊字符
        # \uFF00-\uFFEF 包括全角ASCII、全角标点、半角假名等
        # 这里也包括了破折号——和点号·
        pattern = r'(?<=[\u4e00-\u9fff\u3000-\u303f\uFF00-\uFFEF——·\n])[^\S\n]+(?=[\u4e00-\u9fff\u3000-\u303f\uFF00-\uFFEF——·\n])'
        text = re.sub(pattern, '', text)
        text = re.sub(r'\n+', '\n', text)
        # 去除每行的前后空白字符
        text = '\n'.join(line.strip() for line in text.split('\n'))
        return text
    
    def zh_traditional_to_simplified(self, text):
        # Convert traditional Chinese to simplified Chinese
        converter = opencc.OpenCC('t2s')
        simplified_text = converter.convert(text)
        return simplified_text

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python extract_text.py <file_path>")
        sys.exit(1)

    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        sys.exit(1)

    try:
        parser = DocumentParser(file_path)
        text = parser.extract_text()
        print(f"Extracted text from {file_path}:\n")
        print(text)
    except ValueError as e:
        print(e)
        sys.exit(1)
